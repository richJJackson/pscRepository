# -*- coding: utf-8 -*-
"""Insert psc R package slide into PSC_talk_academic.pptx."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor

PATH = "PSC_talk_academic.pptx"
INSERT_AT = 8  # before "Applications and case studies"

ML, MR = 0.55, 0.55
TITLE_T, TITLE_H = 0.28, 0.72
BODY_T = 1.05
LEFT_W = 5.85
RIGHT_L = 6.55
RIGHT_W = 6.25
BODY_H = 6.0

NAVY = RGBColor(0x0E, 0x28, 0x41)
TEAL = RGBColor(0x15, 0x60, 0x82)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
CODE_BG = RGBColor(0xF0, 0xF4, 0xF8)


def inch(v):
    return int(Inches(v))


def move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])


def style_title(shape):
    tf = shape.text_frame
    for p in tf.paragraphs:
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = NAVY


def fill_bullets(tf, lines, size=Pt(17)):
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = size
        p.font.color.rgb = DARK_GREY
        p.space_after = Pt(8)


def fill_code(tf, code, size=Pt(13)):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.name = "Menlo"
    p.font.size = size
    p.font.color.rgb = DARK_GREY


def add_title_accent(slide):
    bar = slide.shapes.add_shape(1, inch(ML), inch(TITLE_T + TITLE_H + 0.05), inch(2.2), inch(0.06))
    bar.name = "TitleAccent"
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xE9, 0x71, 0x32)
    bar.line.fill.background()


def update_outline(slide):
    for shape in slide.placeholders:
        if not shape.has_text_frame or shape == slide.shapes.title:
            continue
        tf = shape.text_frame
        lines = [p.text for p in tf.paragraphs if p.text.strip()]
        if not any("psc R package" in ln.lower() for ln in lines):
            p = tf.add_paragraph()
            p.text = "The psc R package (software)"
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GREY
            # move before applications line
            paras = list(tf.paragraphs)
            for i, para in enumerate(paras):
                if "Applications in practice" in para.text:
                    # rebuild order: insert software before applications
                    all_text = [pp.text for pp in paras if pp.text.strip()]
                    if "The psc R package (software)" not in all_text:
                        idx = all_text.index("Applications in practice")
                        all_text.insert(idx, "The psc R package (software)")
                        tf.clear()
                        for j, t in enumerate(all_text):
                            pp = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                            pp.text = t
                            pp.font.size = Pt(18)
                            pp.font.color.rgb = DARK_GREY
                            pp.space_after = Pt(10)
        break


def main():
    prs = Presentation(PATH)

    # skip if already added
    for slide in prs.slides:
        if slide.shapes.title and "psc R package" in slide.shapes.title.text.lower():
            print("Software slide already present; updating outline only.")
            update_outline(prs.slides[1])
            prs.save(PATH)
            return

    layout = prs.slide_layouts[1]  # Title and Content
    prs.slides.add_slide(layout)
    move_slide(prs, len(prs.slides) - 1, INSERT_AT)
    slide = prs.slides[INSERT_AT]

    slide.shapes.title.text = "The psc R package"
    style_title(slide.shapes.title)
    slide.shapes.title.left = inch(ML)
    slide.shapes.title.top = inch(TITLE_T)
    slide.shapes.title.width = inch(13.33 - ML - MR)
    slide.shapes.title.height = inch(TITLE_H)
    add_title_accent(slide)

    # remove default body placeholder content area - use custom boxes
    for ph in list(slide.placeholders):
        if ph.has_text_frame and ph != slide.shapes.title:
            ph.text_frame.text = ""

    bullets = [
        "R package for Personalised Synthetic Controls (v2.0.1)",
        "Compare a data cohort to a counterfactual model (CFM)",
        "Supports flexsurvreg survival models and glm outcomes",
        "Bayesian MCMC for treatment-effect parameter beta",
        "Key functions: pscCFM(), pscfit(), summary(), plot()",
        "Related: pscDesign (trial simulation) in development",
        "Install: devtools::install_github(\"richJJackson/psc\")",
        "Docs: https://richjjackson.github.io/psc/",
    ]

    left = slide.shapes.add_textbox(inch(ML), inch(BODY_T), inch(LEFT_W), inch(BODY_H))
    fill_bullets(left.text_frame, bullets)

    code = (
        "library(psc)\n"
        "library(survival)\n"
        "\n"
        "# counterfactual model (from control data)\n"
        "CFM <- pscCFM(surv.mod)\n"
        "\n"
        "# evaluation cohort (e.g. single arm)\n"
        "fit <- pscfit(CFM, data)\n"
        "\n"
        "summary(fit)\n"
        "plot(fit)"
    )

    panel = slide.shapes.add_shape(
        1, inch(RIGHT_L), inch(BODY_T), inch(RIGHT_W), inch(BODY_H)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = CODE_BG
    panel.line.color.rgb = TEAL

    code_box = slide.shapes.add_textbox(
        inch(RIGHT_L + 0.15), inch(BODY_T + 0.15), inch(RIGHT_W - 0.3), inch(BODY_H - 0.3)
    )
    fill_code(code_box.text_frame, code)

    cap = slide.shapes.add_textbox(
        inch(RIGHT_L), inch(BODY_T + BODY_H - 0.35), inch(RIGHT_W), inch(0.3)
    )
    cap.text_frame.text = "Minimal workflow"
    for p in cap.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = TEAL
        p.alignment = PP_ALIGN.CENTER

    update_outline(prs.slides[1])

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "psc package on GitHub richJJackson/psc. "
        "pscCFM wraps a fitted model for sharing; pscfit runs the comparison. "
        "Mention vignette and gemCFM example data in package."
    )

    prs.save(PATH)
    print(f"Added psc software slide at position {INSERT_AT + 1} ({len(prs.slides)} slides total)")


if __name__ == "__main__":
    main()
