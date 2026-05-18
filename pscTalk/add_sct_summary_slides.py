# -*- coding: utf-8 -*-
"""Add 4 SCT summary slides from the manuscript Word document."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DST = "PSC_talk_academic.pptx"
REPLACE_START = 23  # 0-based index of first SCT slide (slide 24)
NUM_OLD = 2

ML, MR = 0.55, 0.55
TITLE_T, TITLE_H = 0.28, 0.72
BODY_T = 1.05
BODY_H = 5.85

NAVY = RGBColor(0x0E, 0x28, 0x41)
TEAL = RGBColor(0x15, 0x60, 0x82)
DARK = RGBColor(0x33, 0x33, 0x33)
ORANGE = RGBColor(0xE9, 0x71, 0x32)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HDR_BG = RGBColor(0x15, 0x60, 0x82)


def inch(v):
    return int(Inches(v))


def delete_slide(prs, index):
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)
    slide_id = slides[index]
    r_id = slide_id.rId
    prs.part.drop_rel(r_id)
    slide_id_list.remove(slide_id)


def move_slide(prs, old_index, new_index):
    xml = prs.slides._sldIdLst
    slides = list(xml)
    xml.remove(slides[old_index])
    xml.insert(new_index, slides[old_index])


def style_title(shape, text):
    shape.text = text
    for p in shape.text_frame.paragraphs:
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = NAVY


def add_accent(slide):
    bar = slide.shapes.add_shape(1, inch(ML), inch(TITLE_T + TITLE_H + 0.04), inch(2.2), inch(0.06))
    bar.name = "TitleAccent"
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()


def new_slide(prs, title, insert_at):
    layout = prs.slide_layouts[1]
    prs.slides.add_slide(layout)
    move_slide(prs, len(prs.slides) - 1, insert_at)
    slide = prs.slides[insert_at]
    style_title(slide.shapes.title, title)
    slide.shapes.title.left = inch(ML)
    slide.shapes.title.top = inch(TITLE_T)
    slide.shapes.title.width = inch(13.33 - ML - MR)
    slide.shapes.title.height = inch(TITLE_H)
    add_accent(slide)
    for ph in slide.placeholders:
        if ph.has_text_frame and ph != slide.shapes.title:
            ph.text_frame.text = ""
    return slide


def add_bullets(slide, lines, top=BODY_T, compact=False):
    box = slide.shapes.add_textbox(inch(ML), inch(top), inch(13.33 - ML - MR), inch(BODY_H - (top - BODY_T)))
    tf = box.text_frame
    tf.word_wrap = True
    size = Pt(17) if compact else Pt(18)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.color.rgb = DARK
        p.space_after = Pt(8)
    return box


def add_simulation_table(slide):
    """Simplified sample-size table (80% power, alpha=0.05) from manuscript."""
    rows = [
        ["Design", "Total N (control)", "Notes"],
        ["Fully synthetic", "63 (52)", "Single arm vs external synthetic control"],
        ["Partially synthetic 1:1", "160 (107)", "Smallest partial design; validation arm"],
        ["Partially synthetic 1:2", "152 (99)", "Intermediate allocation"],
        ["Partially synthetic 1:3", "147 (94)", "More experimental patients"],
        ["Standard RCT", "198 (95)", "Concurrent randomised control"],
    ]
    cols, nrows = 3, len(rows)
    left, top, width, height = inch(ML), inch(1.55), inch(12.2), inch(3.2)
    table = slide.shapes.add_table(nrows, cols, left, top, width, height).table

    col_widths = [Inches(4.2), Inches(2.8), Inches(5.2)]
    for j, w in enumerate(col_widths):
        table.columns[j].width = w

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13 if i > 0 else 14)
                p.font.bold = i == 0
                p.font.color.rgb = WHITE if i == 0 else DARK
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = HDR_BG

    note = slide.shapes.add_textbox(inch(ML), inch(4.85), inch(12.2), inch(2.0))
    tf = note.text_frame
    tf.word_wrap = True
    lines = [
        "Simulation: two cohorts (n=5,000 each) with deliberate covariate imbalance between external and prospective populations",
        "Compared population-averaged SC, personalised SC (PSC), Bayesian hybrid, and unadjusted estimators",
        "Under misspecified exchangeability, synthetic control methods outperformed unadjusted comparisons (see manuscript)",
        "Figures in parentheses: total sample (patients on control) at 80% power, one-sided alpha=0.05",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)


def update_outline(prs):
    lines = [
        "The problem: evidence without a concurrent control arm",
        "Synthetic controls - recap",
        "Personalised synthetic controls (PSC)",
        "Patient-level prediction and likelihood",
        "Bayesian estimation (MCMC)",
        "PSC as causal inference",
        "Comparison with other causal / comparative methods",
        "The psc R package (software)",
        "Applications in practice",
        "Synthetically controlled trials (4 slides)",
        "Limitations and future directions",
    ]
    half = (len(lines) + 1) // 2
    s2 = prs.slides[1]
    for sh in s2.shapes:
        if sh.name == "Content Placeholder 2":
            tf = sh.text_frame
            tf.clear()
            for i, t in enumerate(lines[:half]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = t
                p.font.size = Pt(17)
                p.font.color.rgb = DARK
        elif sh.name == "TextBox OutlineRight":
            tf = sh.text_frame
            tf.clear()
            for i, t in enumerate(lines[half:]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = t
                p.font.size = Pt(17)
                p.font.color.rgb = DARK


def already_done(prs):
    n = sum(
        1
        for slide in prs.slides
        if slide.shapes.title and (slide.shapes.title.text or "").startswith("SCT:")
    )
    return n >= 3


def remove_sct_slides(prs):
    """Remove legacy SCT slides by title (end indices first)."""
    targets = {
        "Synthetically controlled trials",
        "SCT design example: unresectable HCC",
    }
    idx = len(prs.slides) - 1
    while idx >= 0:
        t = prs.slides[idx].shapes.title.text if prs.slides[idx].shapes.title else ""
        if t in targets or t.startswith("SCT:") or t.startswith("SCT in practice"):
            delete_slide(prs, idx)
        idx -= 1


def main():
    prs = Presentation(DST)
    if already_done(prs):
        print("SCT summary slides already present.")
        return

    remove_sct_slides(prs)
    insert_at = min(REPLACE_START, len(prs.slides))

    # Insert in reverse order at same index
    slides_data = [
        (
            "SCT in practice: unresectable HCC",
            "bullets",
            [
                "Setting: unresectable HCC; standard of care atezolizumab + bevacizumab",
                "CFM from observational cohort (604 patients); flexible parametric OS and PFS models",
                "Proposed trial: proton radiotherapy at 3 UK centres vs synthetic control",
                "Primary endpoint PFS; target hazard ratio 0.6; one-sided alpha (Phase II)",
                "Design explored with pscDesign: fully synthetic, partially synthetic (1:1, 1:2, 1:3), standard RCT",
                "Staggered recruitment: 10 sites, 18 patients/site/month, 12-month follow-up",
            ],
            "Worked example from manuscript; sample sizes in previous slide.",
        ),
        (
            "SCT: simulation study",
            "table",
            None,
            "Deliberate violation of exchangeability between external and trial populations.",
        ),
        (
            "SCT: design and estimand",
            "bullets",
            [
                "Collect Y(1) prospectively; derive Y(0) from pre-specified external / model-based control",
                "Population-averaged SC: weighted average of external controls (e.g. entropy balance)",
                "Personalised SC (PSC): parametric CFM gives patient-specific counterfactuals",
                "Estimand: ATT in the enrolled cohort, conditional on measured covariates",
                "Assumptions: exchangeability given X, SUTVA, overlap between trial and external data",
                "Design options: fully synthetic (single arm), partially synthetic (+ validation RCT arm), standard RCT",
            ],
            None,
        ),
        (
            "SCT: overview",
            "bullets",
            [
                "Synthetically controlled trials (SCTs) use external evidence as the control arm in prospective studies",
                "Motivation: scarce patients in rare disease and early-phase oncology; prior control information exists",
                "Contrasts with enforced ignorance: historical data and registries increasingly inform design",
                "Builds on synthetic control causal inference (Abadie) and personalised SC (Jackson et al.)",
                "Regulatory and HTA interest in external controls, digital twins, and single-arm enrichment",
                "Goal: transparent prospective design with pre-specified synthetic control before first patient",
            ],
            "From: Synthetically Controlled Trials in Early Phase Clinical Trials (manuscript).",
        ),
    ]

    for title, kind, bullets, notes in slides_data:
        slide = new_slide(prs, title, insert_at)
        if kind == "table":
            add_simulation_table(slide)
        else:
            add_bullets(slide, bullets)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    update_outline(prs)
    prs.save(DST)
    print(f"Added 4 SCT summary slides at position {insert_at + 1}. Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
