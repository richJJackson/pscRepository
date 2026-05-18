# -*- coding: utf-8 -*-
"""Insert Pfizer estimation slides (11-14) and Bayesian MCMC slides into academic deck."""

import os
import tempfile
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

SRC = "Pfizer.pptx"
DST = "PSC_talk_academic.pptx"
INSERT_AT = 6  # after "PSC - methods"

ML, MR = 0.55, 0.55
TITLE_T, TITLE_H = 0.28, 0.72
BODY_T = 1.05
BODY_H = 5.85
LEFT_W = 5.85
RIGHT_L = 6.55
RIGHT_W = 6.25

NAVY = RGBColor(0x0E, 0x28, 0x41)
TEAL = RGBColor(0x15, 0x60, 0x82)
DARK = RGBColor(0x33, 0x33, 0x33)
ORANGE = RGBColor(0xE9, 0x71, 0x32)


def inch(v):
    return int(Inches(v))


def move_slide(prs, old_index, new_index):
    xml = prs.slides._sldIdLst
    slides = list(xml)
    xml.remove(slides[old_index])
    xml.insert(new_index, slides[old_index])


def style_title(shape, text=None):
    if text is not None:
        shape.text = text
    tf = shape.text_frame
    for p in tf.paragraphs:
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = NAVY


def add_accent(slide):
    if any(s.name == "TitleAccent" for s in slide.shapes):
        return
    bar = slide.shapes.add_shape(1, inch(ML), inch(TITLE_T + TITLE_H + 0.04), inch(2.2), inch(0.06))
    bar.name = "TitleAccent"
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()


def set_title_pos(shape):
    shape.left = inch(ML)
    shape.top = inch(TITLE_T)
    shape.width = inch(13.33 - ML - MR)
    shape.height = inch(TITLE_H)


def fill_body(slide, lines, compact=False):
    body = slide.shapes.add_textbox(inch(ML), inch(BODY_T), inch(13.33 - ML - MR), inch(BODY_H))
    tf = body.text_frame
    tf.word_wrap = True
    size = Pt(17) if compact else Pt(18)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.color.rgb = DARK
        p.space_after = Pt(9)
    return body


def add_picture_from_shape(slide, shape, left, top, width, height):
    blob = shape.image.blob
    ext = shape.image.ext or "png"
    path = tempfile.mktemp(suffix=f".{ext}")
    with open(path, "wb") as f:
        f.write(blob)
    try:
        pic = slide.shapes.add_picture(path, inch(left), inch(top), inch(width), inch(height))
    finally:
        try:
            os.remove(path)
        except OSError:
            pass
    return pic


def new_content_slide(prs, title):
    layout = prs.slide_layouts[1]
    prs.slides.add_slide(layout)
    move_slide(prs, len(prs.slides) - 1, INSERT_AT)
    slide = prs.slides[INSERT_AT]
    style_title(slide.shapes.title, title)
    set_title_pos(slide.shapes.title)
    add_accent(slide)
    for ph in slide.placeholders:
        if ph.has_text_frame and ph != slide.shapes.title:
            ph.text_frame.text = ""
    return slide


def copy_pfizer_pictures(pfizer_slide, dest_slide, layout="dual"):
    pics = [s for s in pfizer_slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if not pics:
        return
    if layout == "dual" and len(pics) >= 2:
        add_picture_from_shape(dest_slide, pics[0], ML, 1.15, 5.85, 5.5)
        add_picture_from_shape(dest_slide, pics[1], RIGHT_L, 1.15, RIGHT_W, 5.5)
    elif layout == "single":
        add_picture_from_shape(dest_slide, pics[0], ML + 0.4, 1.15, 12.2, 5.5)
    elif layout == "single_right" and len(pics) >= 1:
        add_picture_from_shape(dest_slide, pics[0], RIGHT_L, 1.15, RIGHT_W, 5.5)


def update_outline(prs):
    lines = [
        "The problem: evidence without a concurrent control arm",
        "Synthetic controls - recap",
        "Personalised synthetic controls (PSC)",
        "PSC as causal inference",
        "Comparison with other causal / comparative methods",
        "Patient-level prediction and likelihood",
        "Bayesian estimation (MCMC)",
        "The psc R package (software)",
        "Applications in practice",
        "Synthetically controlled trials (prospective design)",
        "Limitations and future directions",
    ]
    half = (len(lines) + 1) // 2
    s2 = prs.slides[1]
    for sh in s2.shapes:
        if sh.name == "Content Placeholder 2":
            tf = sh.text_frame
            tf.clear()
            for i, t in enumerate(lines[:half]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = t
                p.font.size = Pt(17)
                p.font.color.rgb = DARK
                p.space_after = Pt(8)
        elif sh.name == "TextBox OutlineRight":
            tf = sh.text_frame
            tf.clear()
            for i, t in enumerate(lines[half:]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = t
                p.font.size = Pt(17)
                p.font.color.rgb = DARK
                p.space_after = Pt(8)


def already_imported(prs):
    for slide in prs.slides:
        if slide.shapes.title and "Patient-level prediction" in (slide.shapes.title.text or ""):
            return True
    return False


def main():
    pfizer = Presentation(SRC)
    prs = Presentation(DST)

    if already_imported(prs):
        print("Estimation slides appear present; refreshing outline only.")
        update_outline(prs)
        prs.save(DST)
        return

    # Insert in reverse order so slides end up 6..10 in correct sequence
    # (each insert at INSERT_AT pushes previous inserts down)

    s = new_content_slide(prs, "Bayesian MCMC procedure")
    fill_body(
        s,
        [
            "Target: posterior P(beta | B, D) proportional to L(D | B, beta) pi(B) pi(beta)",
            "pi(B): MVN from CFM coefficient estimates and variance-covariance matrix",
            "pi(beta): weakly informative (e.g. N(0, 1000) on log hazard ratio)",
            "Each MCMC iteration: draw CFM parameters B; propose beta; compute likelihood ratio",
            "Accept/reject proposal; repeat for nsim iterations (multiple chains, burn-in, thin)",
            "Output: posterior draws for beta — summary() and plot() in the psc package",
        ],
        compact=True,
    )
    s.notes_slide.notes_text_frame.text = (
        "pscfit workflow: pscData -> init -> pscEst -> postSummary. "
        "See ?pscEst and ?pscfit in the psc package."
    )

    s = new_content_slide(prs, "Why Bayesian estimation?")
    fill_body(
        s,
        [
            "A plug-in MLE for beta (e.g. optim in R) can appear unbiased in simulation",
            "But the counterfactual model is not known with certainty",
            "CFM coefficients have their own uncertainty and correlation",
            "We need that uncertainty propagated into inference on beta",
            "Bayesian MCMC: informative priors on CFM parameters + weak prior on beta",
        ],
    )

    s = new_content_slide(prs, "Estimation: defining the likelihood")
    fill_body(
        s,
        [
            "Standard survival likelihood, conditional on CFM parameters and patient covariates",
            "For each patient i: contribution from event time and censoring indicator",
            "Linear predictor: Gamma_i = X_i' gamma + beta",
            "gamma from the counterfactual model; beta is the treatment-effect parameter",
            "Product over patients: L(D | Lambda, Gamma) compared to CFM predictions",
        ],
    )

    s = new_content_slide(prs, "Patient-level prediction (example)")
    copy_pfizer_pictures(pfizer.slides[11], s, layout="single")

    s = new_content_slide(prs, "Patient-level prediction")
    copy_pfizer_pictures(pfizer.slides[10], s, layout="dual")
    cap = s.shapes.add_textbox(inch(ML), inch(0.98), inch(12.2), inch(0.35))
    cap.text_frame.text = (
        "Counterfactual survival from the CFM vs observed outcome for individual patients"
    )
    for p in cap.text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = TEAL

    update_outline(prs)
    prs.save(DST)
    print(f"Inserted Pfizer estimation content at slide {INSERT_AT + 1}+ ({len(prs.slides)} slides total)")


if __name__ == "__main__":
    main()
