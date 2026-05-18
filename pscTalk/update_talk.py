# -*- coding: utf-8 -*-
"""Update PSC_talk_academic.pptx from seminar resources."""

from pptx import Presentation
from pptx.util import Inches, Pt

PATH = "PSC_talk_academic.pptx"
BODY_SIZE = Pt(20)


def move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])


def insert_slide(prs, index, layout_idx=1):
    layout = prs.slide_layouts[layout_idx]
    prs.slides.add_slide(layout)
    move_slide(prs, len(prs.slides) - 1, index)
    return prs.slides[index]


def set_title(slide, text):
    if slide.shapes.title:
        slide.shapes.title.text = text


def find_body_placeholder(slide):
    title = slide.shapes.title
    for shape in slide.placeholders:
        if shape.has_text_frame and shape != title:
            return shape
    for shape in slide.shapes:
        if shape.has_text_frame and shape != title:
            return shape
    return None


def set_body(slide, lines, font_size=BODY_SIZE):
    ph = find_body_placeholder(slide)
    if ph is None:
        return False
    tf = ph.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = font_size
    return True


def add_body_textbox(slide, lines, left=Inches(0.6), top=Inches(1.6), width=Inches(8.5), height=Inches(5.0)):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = BODY_SIZE
    return box


def set_notes(slide, text):
    if not text.strip():
        return
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def main():
    prs = Presentation(PATH)

    s1 = prs.slides[0]
    set_title(s1, "Personalised Synthetic Controls")
    sub = s1.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.5), Inches(1.2))
    sub.text_frame.text = (
        "Causal inference and prospective clinical trial design\n"
        "Dr Richard Jackson"
    )
    for p in sub.text_frame.paragraphs:
        p.font.size = Pt(22)
    set_notes(s1, "45-minute seminar: PSC, causal framing, comparisons, applications, SCT design.")

    s2 = prs.slides[1]
    set_title(s2, "Outline")
    set_body(
        s2,
        [
            "The problem: evidence without a concurrent control arm",
            "Synthetic controls - recap",
            "Personalised synthetic controls (PSC)",
            "PSC as causal inference",
            "Comparison with other causal / comparative methods",
            "Applications in practice",
            "Synthetically controlled trials (prospective design)",
            "Limitations and future directions",
        ],
    )

    s3 = prs.slides[2]
    set_title(s3, "The problem")
    add_body_textbox(
        s3,
        [
            "RCTs remain the gold standard but require large samples, time, and cost",
            "Scarce resources in rare disease and early-phase oncology",
            "Single-arm trials are common but lack a concurrent comparator for efficacy",
            "Substantial prior information on control care exists (historical trials, registries)",
            "Challenge: credible causal estimate of experimental treatment vs standard of care",
        ],
    )
    set_notes(s3, "From SCT manuscript: enforced ignorance vs information available at design.")

    s4 = prs.slides[3]
    set_title(s4, "Synthetic controls - a recap")
    set_body(
        s4,
        [
            "Causal inference tool (Abadie): construct a synthetic control from donor units",
            "Donors weighted so treated and synthetic cohorts match pre-treatment covariates",
            "Typically targets a population-averaged treatment effect (ATE)",
            "Increasing use in healthcare and cancer research (Thorlund et al.)",
            "Limitations: few donors, heterogeneous patients, unstable weights under poor overlap",
        ],
    )

    s5 = prs.slides[4]
    set_title(s5, "Personalised synthetic controls")
    set_body(
        s5,
        [
            "Parametric counterfactual model (CFM) from external / historical control data",
            "Patient-level predicted outcome under control given covariates X",
            "Compare observed outcomes in evaluation cohort to CFM predictions",
            "Beta: systematic deviation from counterfactual (treatment / policy effect)",
            "Personalised linear predictor (Jackson et al., BMC Med Res Methodol 2025)",
        ],
    )

    s6 = prs.slides[5]
    set_title(s6, "PSC - methods")
    set_body(
        s6,
        [
            "Likelihood: product of patient-level survival or GLM contributions",
            "Bayesian estimation of beta; uncertainty in CFM parameters incorporated",
            "CFM often from RCT control arm (flexible parametric survival, 11+ covariates)",
            "Subgroup and transportability analyses via row indices / weights",
            "Shareable CFM object without releasing individual development data",
        ],
    )

    s_causal = insert_slide(prs, 6, layout_idx=1)
    set_title(s_causal, "PSC as causal inference")
    set_body(
        s_causal,
        [
            "Potential outcomes Y(0) and Y(1); observe Y(1), estimate Y(0) via CFM",
            "Estimand: ATT conditional on measured covariates X",
            "Identification: exchangeability given X, SUTVA, covariate overlap / transportability",
            "Counterfactual control fixed before enrolment; depends on realised trial population",
            "Beta on log-hazard scale: contrast vs counterfactual predictions",
        ],
    )
    set_notes(s_causal, "Pfizer deck: Rubin framework; HR as contrast of potential outcomes.")

    s_compare = insert_slide(prs, 7, layout_idx=1)
    set_title(s_compare, "How PSC compares with other methods")
    set_body(
        s_compare,
        [
            "RCT: randomisation gives ITT treatment effect (gold standard when feasible)",
            "Propensity scores / IPW: balance observational cohorts to ATE or ATT",
            "MAIC / STC: align two trial populations for indirect comparison A vs B",
            "Classical synthetic controls: donor weights for population-level ATE",
            "PSC: validated CFM + evaluation cohort gives beta vs control model (not A vs B)",
        ],
    )
    set_notes(s_compare, "HDSforum comparison; different estimand from MAIC/STC.")

    s_apps = prs.slides[8]
    set_title(s_apps, "Applications and case studies")
    set_body(
        s_apps,
        [
            "Retrospective: single-arm or one-arm vs external control model",
            "Calibration: PSC on control arm of RCT (ESPAC Gem arm)",
            "Benchmarking: observational cohort vs known trial (HCC Atezo-Bev)",
            "Prediction: Phase II single-arm HR vs Phase III RCT",
        ],
    )

    s_espac1 = prs.slides[9]
    set_title(s_espac1, "Application: ESPAC trials")
    set_body(
        s_espac1,
        [
            "Step 1 - Study population: are Gem development data and ESPAC-4 comparable?",
            "Step 2 - Model validation: PSC on Gem arm gives calibration HR approx 0.98",
            "Step 3 - GemCap vs Gem: PSC HR comparable to direct RCT comparison",
            "Greater uncertainty under PSC despite similar effective sample size",
        ],
    )

    s_sct = prs.slides[17]
    set_title(s_sct, "Synthetically controlled trials")
    set_body(
        s_sct,
        [
            "Prospective trial: collect Y(1); derive Y(0) from pre-specified synthetic control",
            "Fully synthetic (single arm), partially synthetic (validation arm), or hybrid",
            "Population-averaged SC (entropy weights) vs personalised SC (CFM)",
            "Design-stage control over eligibility to ensure covariate overlap",
            "Growing regulatory and HTA interest in external / synthetic controls",
        ],
    )

    s_design = prs.slides[18]
    set_title(s_design, "SCT design example: unresectable HCC")
    set_body(
        s_design,
        [
            "SOC: atezolizumab + bevacizumab; CFM from 604-patient observational cohort",
            "Evaluate proton radiotherapy at 3 UK centres vs synthetic control (PFS primary)",
            "Target HR = 0.6; one-sided alpha typical of Phase II",
            "pscDesign: fully RCT vs partially synthetic (1:1, 1:2, 1:3) vs fully synthetic",
            "Staggered recruitment: 10 sites, 18 patients/site/month, 12-month follow-up",
        ],
    )
    set_notes(s_design, "SCT manuscript worked example; 80% and 90% power scenarios.")

    s_lim = prs.slides[19]
    set_title(s_lim, "Limitations")
    set_body(
        s_lim,
        [
            "Requires overlap: unstable weights (population SC) or CFM extrapolation (PSC)",
            "Restrict eligibility at design to regions supported by the synthetic control",
            "Unmeasured confounding; outcome definitions and SOC context must align",
            "SCT evidence is not equivalent to a fully randomised trial; efficiency has a cost",
            "Model validation and transport diagnostics should be pre-specified",
        ],
    )

    s_future = prs.slides[20]
    set_title(s_future, "Further work")
    set_body(
        s_future,
        [
            "Proximity-weighted PSC for covariate shift and transportability",
            "Simulation under deliberate exchangeability violations (pscDesign package)",
            "Multiple treatments, non-proportional hazards, alternative effect measures",
            "Small randomised control subsample for prospective validation",
            "Software: psc R package (Jackson et al., BMC Med Res Methodol 2025;25:91)",
        ],
    )

    prs.save(PATH)
    print(f"Updated {PATH} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
