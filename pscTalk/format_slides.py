# -*- coding: utf-8 -*-
"""Reformat PSC_talk_academic.pptx: fix text/figure overlap and improve styling."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

PATH = "PSC_talk_academic.pptx"

# Theme-aligned colours (from theme1.xml)
NAVY = RGBColor(0x0E, 0x28, 0x41)
TEAL = RGBColor(0x15, 0x60, 0x82)
ORANGE = RGBColor(0xE9, 0x71, 0x32)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)

TITLE_PT = Pt(32)
BODY_PT = Pt(18)
BODY_COMPACT_PT = Pt(16)
CAPTION_PT = Pt(14)

# Layout (inches)  13.33" x 7.5" slide
ML = 0.55
MR = 0.55
TITLE_T = 0.28
TITLE_H = 0.72
BODY_T = 1.05
BODY_H = 6.15
LEFT_W = 5.75
RIGHT_L = 6.55
RIGHT_W = 6.25
FIG_T = 1.05
FIG_H = 6.0


def inch(v):
    return int(Inches(v))


def set_pos(shape, left, top, width, height):
    shape.left = inch(left)
    shape.top = inch(top)
    shape.width = inch(width)
    shape.height = inch(height)


def style_title(shape):
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.font.size = TITLE_PT
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.space_after = Pt(4)
    tf.margin_left = tf.margin_right = Pt(0)


def style_body_frame(tf, size=BODY_PT, color=DARK_GREY):
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.size = size
            run.font.color.rgb = color
        if not p.runs and p.text:
            p.font.size = size
            p.font.color.rgb = color
        p.space_after = Pt(10)
        p.line_spacing = 1.15
        p.level = min(p.level, 1)


def style_body_shape(shape, compact=False):
    if not shape or not shape.has_text_frame:
        return
    style_body_frame(shape.text_frame, BODY_COMPACT_PT if compact else BODY_PT)


def add_text_panel(slide, left, top, width, height):
    """Light panel behind text for readability on busy slides."""
    panel = slide.shapes.add_shape(
        1, inch(left), inch(top), inch(width), inch(height)  # MSO_SHAPE.RECTANGLE = 1
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = LIGHT_BG
    panel.line.fill.background()
    # send to back
    sp_tree = slide.shapes._spTree
    sp_tree.remove(panel._element)
    sp_tree.insert(2, panel._element)
    return panel


def get_title_shape(slide):
    if slide.shapes.title:
        return slide.shapes.title
    for s in slide.shapes:
        if s.name == "Title 1" and s.has_text_frame:
            return s
    return None


def get_body_placeholder(slide):
    title = get_title_shape(slide)
    for s in slide.placeholders:
        if s.has_text_frame and s != title:
            return s
    return None


def add_title_accent(slide):
    """Thin accent line under slide title."""
    for sh in slide.shapes:
        if sh.name == "TitleAccent":
            return
    bar = slide.shapes.add_shape(
        1, inch(ML), inch(TITLE_T + TITLE_H + 0.05), inch(2.2), inch(0.06)
    )
    bar.name = "TitleAccent"
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()


def format_text_only_slide(slide, compact=False):
    t = get_title_shape(slide)
    if t:
        set_pos(t, ML, TITLE_T, 13.33 - ML - MR, TITLE_H)
        style_title(t)
        add_title_accent(slide)
    body = get_body_placeholder(slide)
    if body:
        set_pos(body, ML, BODY_T, 13.33 - ML - MR, BODY_H)
        style_body_shape(body, compact=compact)


def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def main():
    prs = Presentation(PATH)
    slides = prs.slides

    # --- Slide 1: title ---
    s1 = slides[0]
    t = get_title_shape(s1)
    style_title(t)
    # remove duplicate textbox; use subtitle placeholder
    for sh in list(s1.shapes):
        if sh.name == "TextBox 3":
            remove_shape(sh)
    for sh in s1.placeholders:
        if sh.name == "Subtitle 2" and sh.has_text_frame:
            sh.text_frame.text = (
                "Causal inference and prospective clinical trial design\n"
                "Dr Richard Jackson"
            )
            for p in sh.text_frame.paragraphs:
                p.font.size = Pt(22)
                p.font.color.rgb = TEAL
                p.alignment = PP_ALIGN.CENTER
            set_pos(sh, ML, 1.35, 13.33 - 2 * ML, 1.4)

    # --- Slides 2, 4-9, 18-21: text-only ---
    text_only = [1, 3, 4, 5, 6, 7, 8, 17, 18, 19, 20]
    compact_slides = {1, 7, 8, 19, 20}  # many bullets
    for idx in text_only:
        format_text_only_slide(slides[idx], compact=(idx in compact_slides))

    # --- Slide 3: problem  left text, right figure ---
    s3 = slides[2]
    style_title(get_title_shape(s3))
    for sh in s3.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            set_pos(sh, RIGHT_L, FIG_T, RIGHT_W, FIG_H)
        elif sh.name == "TextBox 4":
            set_pos(sh, ML, BODY_T, LEFT_W, BODY_H)
            style_body_shape(sh)

    # --- Slide 10: ESPAC  left bullets, right figure ---
    s10 = slides[9]
    style_title(get_title_shape(s10))
    for sh in s10.shapes:
        if sh.name == "TextBox 5":
            set_pos(sh, ML, BODY_T, LEFT_W, BODY_H - 0.2)
            style_body_shape(sh, compact=True)
        elif sh.name == "Content Placeholder 4":
            set_pos(sh, RIGHT_L, FIG_T, RIGHT_W, FIG_H)

    # --- Slide 11: ESPAC step 3  left text, right chart ---
    s11 = slides[10]
    style_title(get_title_shape(s11))
    for sh in s11.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            set_pos(sh, RIGHT_L, FIG_T, RIGHT_W, FIG_H)
        elif sh.name == "Content Placeholder 2":
            set_pos(sh, ML, BODY_T, LEFT_W, 0.55)
            style_body_frame(sh.text_frame, Pt(20), TEAL)
            for p in sh.text_frame.paragraphs:
                p.font.bold = True
        elif sh.name == "TextBox 6":
            set_pos(sh, ML, 1.75, LEFT_W, 5.2)
            style_body_shape(sh, compact=True)

    # --- Slide 12: HCC  fix duplicate titles, side-by-side figures ---
    s12 = slides[11]
    main_title = get_title_shape(s12)
    for sh in list(s12.shapes):
        if sh.shape_type != MSO_SHAPE_TYPE.TEXT_BOX or not sh.has_text_frame:
            continue
        if sh is main_title:
            continue
        if "Application: HCC" in sh.text or sh.name == "Title 1":
            remove_shape(sh)
    t = get_title_shape(s12)
    if t:
        t.text = "Application: HCC (Atezo-Bev vs Sorafenib)"
        style_title(t)
    pics = [sh for sh in s12.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if len(pics) >= 2:
        set_pos(pics[0], ML, FIG_T, 5.85, FIG_H - 0.3)
        set_pos(pics[1], RIGHT_L, FIG_T, RIGHT_W, FIG_H - 0.3)

    # --- Slide 13: summary strip + figures below ---
    s13 = slides[12]
    style_title(get_title_shape(s13))
    for sh in s13.shapes:
        if sh.name == "Content Placeholder 2" and sh.has_text_frame:
            set_pos(sh, ML, BODY_T, 13.33 - ML - MR, 0.85)
            style_body_frame(sh.text_frame, Pt(17), DARK_GREY)
        elif sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if sh.left < inch(6):
                set_pos(sh, ML, 2.05, 5.9, 4.85)
            else:
                set_pos(sh, RIGHT_L, 1.85, RIGHT_W, 5.2)
        elif sh.name.startswith("Doughnut"):
            set_pos(sh, RIGHT_L, 6.55, RIGHT_W, 0.7)

    # --- Slide 14: IMBRAVE eligibility  title bar + dual figures ---
    s14 = slides[13]
    imb_title = None
    for sh in list(s14.shapes):
        if sh.has_text_frame and "IMBRAVE" in (sh.text or ""):
            if imb_title is None:
                imb_title = sh
                sh.text_frame.text = "HCC: eligibility subgroups (IMBRAVE 150)"
                set_pos(sh, ML, TITLE_T, 13.33 - ML - MR, TITLE_H)
                style_title(sh)
            else:
                remove_shape(sh)
    if imb_title is None:
        tb = s14.shapes.add_textbox(inch(ML), inch(TITLE_T), inch(12), inch(TITLE_H))
        tb.text_frame.text = "HCC: eligibility subgroups (IMBRAVE 150)"
        style_title(tb)
    add_title_accent(s14)
    for sh in s14.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if sh.left < inch(6):
                set_pos(sh, ML, 1.85, 5.95, 5.1)
            else:
                set_pos(sh, RIGHT_L, 1.85, RIGHT_W, 5.1)
        elif sh.name.startswith("TextBox") and "Eligible" in (sh.text or ""):
            set_pos(sh, sh.left / 914400, 1.35, 3.5, 0.45)
            style_body_frame(sh.text_frame, CAPTION_PT, NAVY)
            for p in sh.text_frame.paragraphs:
                p.font.bold = True
        elif sh.name.startswith("Rectangle"):
            set_pos(sh, 0, 0, 13.33, 0.42)
            sh.fill.solid()
            sh.fill.fore_color.rgb = NAVY

    # --- Slides 15-17: Phase II/III  title, labels, then figures ---
    for idx in [14, 15, 16]:
        s = slides[idx]
        style_title(get_title_shape(s))
        for sh in s.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if sh.left < inch(6):
                    set_pos(sh, ML, 1.75, 5.85, 5.35)
                else:
                    set_pos(sh, RIGHT_L, 1.75, RIGHT_W, 5.35)
            elif sh.name.startswith("TextBox"):
                txt = sh.text or ""
                if "Single Arm" in txt or "Phase III" in txt:
                    set_pos(sh, sh.left / 914400, 1.22, 4.5, 0.4)
                    style_body_frame(sh.text_frame, CAPTION_PT, TEAL)
                    for p in sh.text_frame.paragraphs:
                        p.font.bold = True
                elif "CI" in txt or "HR" in txt or "0.63" in txt:
                    set_pos(sh, sh.left / 914400, 1.55, 4.0, 0.35)
                    style_body_frame(sh.text_frame, Pt(15), ORANGE)
                    for p in sh.text_frame.paragraphs:
                        p.font.bold = True
            elif sh.name == "Content Placeholder 3":
                # hide empty placeholder area
                set_pos(sh, ML, 1.75, 5.85, 5.35)

    # --- Slide 22: closing ---
    s22 = slides[21]
    t = get_title_shape(s22)
    if t:
        t.text = "Thank you"
        style_title(t)
    body = get_body_placeholder(s22)
    if body:
        body.text_frame.text = "Questions and discussion"
        set_pos(body, ML, 2.2, 13.33 - 2 * ML, 1.5)
        for p in body.text_frame.paragraphs:
            p.font.size = Pt(24)
            p.font.color.rgb = TEAL
            p.alignment = PP_ALIGN.CENTER

    # Normalise title positions on figure slides
    for idx in [9, 10, 11, 12, 13, 14, 15, 16]:
        t = get_title_shape(slides[idx])
        if t:
            set_pos(t, ML, TITLE_T, 13.33 - ML - MR, TITLE_H)
            style_title(t)

    # --- Second pass: fix residual overlaps ---
    s13 = slides[12]
    for sh in s13.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE and sh.left >= inch(RIGHT_L - 0.1):
            set_pos(sh, RIGHT_L, 2.05, RIGHT_W, 4.85)

    for idx in [14, 15, 16]:
        s = slides[idx]
        fig_top = 2.05
        for sh in s.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if sh.left < inch(6):
                    set_pos(sh, ML, fig_top, 5.85, 5.0)
                else:
                    set_pos(sh, RIGHT_L, fig_top, RIGHT_W, 5.0)
            elif sh.name.startswith("TextBox"):
                txt = sh.text or ""
                if "Single Arm" in txt:
                    set_pos(sh, ML + 0.2, 1.12, 5.4, 0.38)
                elif txt.strip() == "Phase III":
                    set_pos(sh, RIGHT_L + 0.2, 1.12, 5.4, 0.38)
                elif "CI" in txt or "HR" in txt or "0.63" in txt:
                    col_left = ML + 0.2 if sh.left < inch(6) else RIGHT_L + 0.2
                    set_pos(sh, col_left, 1.52, 5.2, 0.38)
                style_body_frame(sh.text_frame, Pt(15), ORANGE if "CI" in txt or "HR" in txt else TEAL)
                for p in sh.text_frame.paragraphs:
                    p.font.bold = True
            elif sh.name == "Content Placeholder 3":
                set_pos(sh, ML, fig_top, 5.85, 5.0)

    # --- Outline: two-column layout for readability ---
    s2 = slides[1]
    body = get_body_placeholder(s2)
    if body and body.has_text_frame:
        lines = [p.text for p in body.text_frame.paragraphs if p.text.strip()]
        half = (len(lines) + 1) // 2
        left_lines = lines[:half]
        right_lines = lines[half:]
        body.text_frame.clear()
        for i, line in enumerate(left_lines):
            p = body.text_frame.paragraphs[0] if i == 0 else body.text_frame.add_paragraph()
            p.text = line
            p.level = 0
        set_pos(body, ML, BODY_T, 6.0, BODY_H)
        style_body_shape(body)
        # remove old right column box if re-running
        for sh in list(s2.shapes):
            if sh.name == "TextBox OutlineRight":
                remove_shape(sh)
        right_box = s2.shapes.add_textbox(inch(6.85), inch(BODY_T), inch(6.0), inch(BODY_H))
        right_box.name = "TextBox OutlineRight"
        tf = right_box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(right_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.level = 0
        style_body_shape(right_box)

    prs.save(PATH)
    print(f"Formatted {PATH}")


if __name__ == "__main__":
    main()
