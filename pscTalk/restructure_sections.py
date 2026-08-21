# -*- coding: utf-8 -*-
"""Restructure Newcastle.pptx into four seminar sections with section dividers."""

import shutil

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

PATH = "pscTalk/Newcastle.pptx"
BACKUP = "pscTalk/Newcastle_backup.pptx"

ML, MR = 0.55, 0.55
TITLE_T, TITLE_H = 0.28, 0.72
BODY_T = 1.05
BODY_H = 6.0
LEFT_W = 5.85
RIGHT_L = 6.55
RIGHT_W = 6.25

BG = RGBColor(0xF8, 0xF8, 0xF5)
TITLE = RGBColor(0x18, 0x21, 0x2E)
TEXT = RGBColor(0x21, 0x29, 0x36)
SUBTLE = RGBColor(0x42, 0x4B, 0x59)
ACCENT = RGBColor(0x15, 0x60, 0x82)
CODE_BG = RGBColor(0xF0, 0xF4, 0xF8)

SECTIONS = [
    ("Foundations", "Theory, Models & Estimation"),
    ("Hands-On", "Your PSC Workflow"),
    ("In Practice", "Applications & Case Studies"),
    ("The pscLibrary", "Open Models, Code & Data"),
]

OUTLINE_LINES = [
    "Foundations — Theory, Models & Estimation",
    "Hands-On — Your PSC Workflow",
    "In Practice — Applications & Case Studies",
    "The pscLibrary — Open Models, Code & Data",
]


def inch(v):
    return int(Inches(v))


def move_slide(prs, old_index, new_index):
    xml = prs.slides._sldIdLst
    slides = list(xml)
    xml.remove(slides[old_index])
    xml.insert(new_index, slides[old_index])


def reorder_slides(prs, order):
    xml = prs.slides._sldIdLst
    slides = list(xml)
    ordered = [slides[i] for i in order]
    for slide in slides:
        xml.remove(slide)
    for slide in ordered:
        xml.append(slide)


def apply_font(p, size_pt, color, bold=False, align=PP_ALIGN.LEFT):
    p.alignment = align
    p.font.name = "Aptos"
    p.font.size = Pt(size_pt)
    p.font.bold = bold
    p.font.color.rgb = color
    for r in p.runs:
        r.font.name = "Aptos"
        r.font.size = Pt(size_pt)
        r.font.bold = bold
        r.font.color.rgb = color


def style_text_frame(tf, size_pt=20, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    tf.word_wrap = True
    for p in tf.paragraphs:
        apply_font(p, size_pt, color, bold=bold, align=align)


def set_slide_background(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_accent(slide):
    for sh in slide.shapes:
        if sh.name == "TitleAccent":
            return
    bar = slide.shapes.add_shape(
        1, inch(ML), inch(TITLE_T + TITLE_H + 0.05), inch(2.2), inch(0.06)
    )
    bar.name = "TitleAccent"
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def clear_body_placeholders(slide):
    title = slide.shapes.title
    for ph in slide.placeholders:
        if ph.has_text_frame and ph != title:
            ph.text_frame.text = ""


def add_bullets(slide, lines, left=ML, top=BODY_T, width=None, size=18):
    if width is None:
        width = 13.33 - ML - MR
    box = slide.shapes.add_textbox(inch(left), inch(top), inch(width), inch(BODY_H))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        apply_font(p, size, TEXT)
        p.space_after = Pt(8)
    return box


def new_content_slide(prs, title_text):
    layout = prs.slide_layouts[1]
    prs.slides.add_slide(layout)
    slide = prs.slides[-1]
    set_slide_background(slide)
    slide.shapes.title.text = title_text
    style_text_frame(slide.shapes.title.text_frame, size_pt=30, color=TITLE, bold=True)
    slide.shapes.title.left = inch(ML)
    slide.shapes.title.top = inch(TITLE_T)
    slide.shapes.title.width = inch(13.33 - ML - MR)
    slide.shapes.title.height = inch(TITLE_H)
    add_title_accent(slide)
    clear_body_placeholders(slide)
    return slide


def add_section_divider(prs, eyebrow, headline):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    prs.slides.add_slide(layout)
    slide = prs.slides[-1]
    set_slide_background(slide)

    for sh in list(slide.shapes):
        if sh.has_text_frame:
            sh.text_frame.text = ""

    eyebrow_box = slide.shapes.add_textbox(inch(ML), inch(2.55), inch(12.2), inch(0.55))
    style_text_frame(eyebrow_box.text_frame, size_pt=22, color=ACCENT, bold=True, align=PP_ALIGN.LEFT)
    eyebrow_box.text_frame.paragraphs[0].text = eyebrow.upper()

    headline_box = slide.shapes.add_textbox(inch(ML), inch(3.15), inch(12.2), inch(1.4))
    style_text_frame(headline_box.text_frame, size_pt=40, color=TITLE, bold=True, align=PP_ALIGN.LEFT)
    headline_box.text_frame.paragraphs[0].text = headline

    bar = slide.shapes.add_shape(1, inch(ML), inch(4.75), inch(2.4), inch(0.07))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    return slide


def add_psc_package_slide(prs):
    slide = new_content_slide(prs, "The psc R package")
    add_bullets(
        slide,
        [
            "R package for Personalised Synthetic Controls",
            "Compare an evaluation cohort to a counterfactual model (CFM)",
            "Supports flexsurvreg survival models and glm outcomes",
            "Bayesian MCMC for treatment-effect parameter beta",
            "Key functions: pscCFM(), pscfit(), summary(), plot()",
            "Related: pscDesign (trial simulation) in development",
            'Install: devtools::install_github("richJJackson/psc")',
            "Docs: https://richjjackson.github.io/psc/",
        ],
        width=LEFT_W,
        size=17,
    )

    panel = slide.shapes.add_shape(1, inch(RIGHT_L), inch(BODY_T), inch(RIGHT_W), inch(BODY_H))
    panel.fill.solid()
    panel.fill.fore_color.rgb = CODE_BG
    panel.line.color.rgb = ACCENT

    code_box = slide.shapes.add_textbox(
        inch(RIGHT_L + 0.15), inch(BODY_T + 0.15), inch(RIGHT_W - 0.3), inch(BODY_H - 0.3)
    )
    code = (
        "library(psc)\n"
        "library(survival)\n\n"
        "CFM <- pscCFM(surv.mod)\n"
        "fit <- pscfit(CFM, data)\n\n"
        "summary(fit)\n"
        "plot(fit)"
    )
    tf = code_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.name = "Menlo"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT
    return slide


def add_workflow_slide(prs):
    slide = new_content_slide(prs, "Five Steps to a PSC Analysis")
    add_bullets(
        slide,
        [
            "1. Specify the counterfactual — fit or load a validated CFM from control data",
            "2. Prepare the evaluation cohort — single-arm trial, registry, or observational data",
            "3. Check compatibility — covariates, scales, and overlap (dataComb / data_match)",
            "4. Estimate the treatment effect — pscfit() with Bayesian inference for beta",
            "5. Report and validate — summary(), plots, and calibration on known controls where possible",
            "",
            "Worked example in this repository: Applications/Gem_vs_GemCap/",
        ],
        size=19,
    )
    return slide


def add_repository_overview_slide(prs):
    slide = new_content_slide(prs, "The pscLibrary Repository")
    add_bullets(
        slide,
        [
            "GitHub: github.com/richJJackson/pscRepository",
            "Companion resource to the psc R package — models, data, and reproducible code",
            "Purpose: share validated CFMs without releasing individual patient data",
            "Enables teaching, benchmarking, and transparent replication of PSC analyses",
            "Each application is self-contained: model object, cohort data, and analysis script",
            "Ideal starting point for adapting PSC to your own disease area or trial",
        ],
        size=19,
    )
    return slide


def add_repository_structure_slide(prs):
    slide = new_content_slide(prs, "What's Inside pscLibrary")
    add_bullets(
        slide,
        [
            "Models/ — pre-validated counterfactual models (e.g. PDAC gemcitabine CFM)",
            "Applications/ — end-to-end worked examples with R scripts and vignettes",
            "Gem vs GemCap (ESPAC) — flagship case study in Applications/Gem_vs_GemCap/",
            "  - flexParaGem.R (CFM), e4_data_cohort.csv, espac_PSC.R, psc_pdac.Rmd",
            "Models include validation scripts; applications include step-by-step workflows",
            "Clone the repo, open the R project (pscLibrary.Rproj), and run the examples",
        ],
        size=18,
    )
    return slide


def update_outline_slide(slide):
    title = slide.shapes.title
    if title:
        title.text = "Talk Outline"
        style_text_frame(title.text_frame, size_pt=30, color=TITLE, bold=True)

    body = None
    for ph in slide.placeholders:
        if ph.has_text_frame and ph != title:
            body = ph
            break

    if body is None:
        body = slide.shapes.add_textbox(inch(ML), inch(BODY_T), inch(13.33 - ML - MR), inch(BODY_H))

    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(OUTLINE_LINES):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        apply_font(p, 20, TEXT)
        p.space_after = Pt(12)

    body.left = inch(ML)
    body.top = inch(BODY_T)
    body.width = inch(13.33 - ML - MR)
    body.height = inch(BODY_H)


def update_references_slide(slide):
    title = slide.shapes.title
    if title:
        title.text = "Key Resources"
        style_text_frame(title.text_frame, size_pt=30, color=TITLE, bold=True)
        add_title_accent(slide)

    refs = [
        "Jackson R, et al. Personalised synthetic controls. BMC Med Res Methodol. 2025;25:91.",
        "Abadie A, et al. Synthetic control methods for comparative case studies. JASA. 2010.",
        "Thorlund K, et al. Synthetic and external controls in clinical trials. BMJ. 2024.",
        "psc R package: https://richjjackson.github.io/psc/",
        "pscLibrary repository: https://github.com/richJJackson/pscRepository",
        "Model portal (PDAC gem): https://richjjackson.github.io/mecPortal/models/pdac_gem.html",
    ]
    add_bullets(slide, refs, size=17)


def apply_uniform_theme(prs):
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    left_margin = int(slide_w * 0.06)
    right_margin = int(slide_w * 0.94)
    title_top = int(slide_h * 0.045)
    title_h = int(slide_h * 0.09)

    for slide in prs.slides:
        set_slide_background(slide)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                if shape.fill.type is not None:
                    shape.fill.background()
            except Exception:
                pass

            tf = shape.text_frame
            is_title = shape == slide.shapes.title or shape.name.lower().startswith("title")
            if is_title and tf.text.strip():
                style_text_frame(tf, size_pt=30, color=TITLE, bold=True)
            elif tf.text.strip():
                style_text_frame(tf, size_pt=18, color=TEXT)


def main():
    shutil.copy2(PATH, BACKUP)

    prs = Presentation(PATH)
    n_orig = len(prs.slides)

    # Append new slides at end: 4 dividers + 4 content slides
    add_section_divider(prs, SECTIONS[0][0], SECTIONS[0][1])
    add_section_divider(prs, SECTIONS[1][0], SECTIONS[1][1])
    add_section_divider(prs, SECTIONS[2][0], SECTIONS[2][1])
    add_section_divider(prs, SECTIONS[3][0], SECTIONS[3][1])
    add_psc_package_slide(prs)
    add_workflow_slide(prs)
    add_repository_overview_slide(prs)
    add_repository_structure_slide(prs)

    idx_div1 = n_orig
    idx_div2 = n_orig + 1
    idx_div3 = n_orig + 2
    idx_div4 = n_orig + 3
    idx_psc = n_orig + 4
    idx_workflow = n_orig + 5
    idx_repo1 = n_orig + 6
    idx_repo2 = n_orig + 7

    # Original 0-based indices (skip empty slide 17 -> index 16)
    methodology = list(range(2, 16)) + list(range(17, 24)) + [45]
    howto_sct = list(range(33, 42))
    examples = [24] + list(range(25, 33)) + list(range(42, 44))
    closing = [44]

    final_order = (
        [0, 1, idx_div1]
        + methodology
        + [idx_div2, idx_psc, idx_workflow]
        + howto_sct
        + [idx_div3]
        + examples
        + [idx_div4, idx_repo1, idx_repo2]
        + closing
    )

    reorder_slides(prs, final_order)

    update_outline_slide(prs.slides[1])

    for slide in prs.slides:
        if slide.shapes.title and slide.shapes.title.text.strip() == "References":
            update_references_slide(slide)
            break

    apply_uniform_theme(prs)
    prs.save(PATH)

    print(f"Restructured {PATH}")
    print(f"Backup saved to {BACKUP}")
    print(f"Slides: {n_orig} -> {len(prs.slides)}")
    print("Sections:")
    for eyebrow, headline in SECTIONS:
        print(f"  - {eyebrow}: {headline}")


if __name__ == "__main__":
    main()
